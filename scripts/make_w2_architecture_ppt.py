"""Generate a layman-friendly Week 2 Architecture Defense PowerPoint."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parents[1] / "W2_Architecture_Defense.pptx"

# Simple clinical / trust palette (not purple AI cliché)
NAVY = RGBColor(0x1B, 0x3A, 0x4B)
TEAL = RGBColor(0x2A, 0x6F, 0x7F)
ACCENT = RGBColor(0xC4, 0x5C, 0x26)
DARK = RGBColor(0x22, 0x22, 0x22)
MUTED = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF4, 0xF7, 0xF8)


def set_run(run, *, size=18, bold=False, color=DARK):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"


def add_bg(slide, color=LIGHT):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def title_bar(slide, text: str):
    shape = slide.shapes.add_shape(
        1,  # rectangle
        Inches(0),
        Inches(0),
        Inches(13.333),
        Inches(1.0),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "  " + text
    set_run(run, size=26, bold=True, color=WHITE)


def bullets(slide, lines: list[str], top=1.3, left=0.7, width=12, size=20):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.8))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(10)
        run = p.add_run()
        run.text = line
        set_run(run, size=size, color=DARK)


def subtitle(slide, text: str, top=1.15):
    box = slide.shapes.add_textbox(Inches(0.7), Inches(top), Inches(12), Inches(0.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    set_run(run, size=16, color=TEAL)


def add_title_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, NAVY)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.5), Inches(2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Clinical Co-Pilot — Week 2"
    set_run(run, size=40, bold=True, color=WHITE)
    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = "Architecture Defense"
    set_run(run2, size=32, bold=False, color=RGBColor(0xB8, 0xD4, 0xDC))
    foot = slide.shapes.add_textbox(Inches(0.8), Inches(5.2), Inches(11), Inches(1))
    tf2 = foot.text_frame
    p = tf2.paragraphs[0]
    run = p.add_run()
    run.text = "Gauntlet AI · AgentForge · Cohort 6\nPlain-language plan: documents + multi-agent + quality gate"
    set_run(run, size=18, color=RGBColor(0xD0, 0xD0, 0xD0))


def add_slide(prs: Presentation, title: str, lines: list[str], note: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT)
    title_bar(slide, title)
    if note:
        subtitle(slide, note)
        bullets(slide, lines, top=1.7, size=19)
    else:
        bullets(slide, lines, top=1.35, size=20)
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)

    add_slide(
        prs,
        "What problem are we solving?",
        [
            "• A doctor has ~90 seconds before walking into the next room.",
            "• The chart has structured data — but the important new info is often in files:",
            "      – a scanned lab PDF",
            "      – a patient intake form from the front desk",
            "• She asks: What changed? What should I watch? What evidence supports that?",
            "• Week 2 makes the AI able to read those documents — carefully, with proof.",
        ],
    )

    add_slide(
        prs,
        "What we already built (Week 1)",
        [
            "• A Clinical Co-Pilot chat next to OpenEMR (the medical record system).",
            "• It answers only from the patient’s chart — every claim has a source.",
            "• It checks permissions (who can see which patient).",
            "• Safety checks (drug interactions, allergies) run in code — not “AI opinion.”",
            "• Week 2 adds on top of this. We do not throw Week 1 away.",
        ],
        note="Baseline we keep",
    )

    add_slide(
        prs,
        "Week 2 in one picture",
        [
            "Doctor asks a question",
            "        ↓",
            "Supervisor (traffic controller)",
            "   ├─ Document reader  →  reads lab PDF / intake form",
            "   ├─ Evidence finder   →  finds clinic guideline snippets",
            "   └─ Chart tools       →  Week 1 patient data",
            "        ↓",
            "Answer with citations  +  quality checks  +  audit trail",
        ],
        note="Small team of specialists — not one giant black box",
    )

    add_slide(
        prs,
        "Part 1 — Reading documents (without making things up)",
        [
            "• Support exactly two file types for now: lab PDF + intake form.",
            "• Save the original file in OpenEMR (the real medical record).",
            "• Use vision AI to fill a strict form (schema) — required fields only.",
            "• If the scan is blurry or a field is unclear → say “couldn’t read confidently.”",
            "• Never invent a lab value or allergy that isn’t on the page.",
            "• Every extracted fact points back to the file (page + quote + highlight box).",
        ],
    )

    add_slide(
        prs,
        "Part 2 — Finding guideline evidence",
        [
            "• We keep a small library of clinic “how we practice” notes",
            "   (diabetes, blood pressure, kidney labs, screening — for our doctor’s patients).",
            "• Search two ways: exact keywords + meaning (hybrid search).",
            "• Re-rank the best snippets; only the top few go to the answer.",
            "• Important rule: patient facts ≠ guideline advice.",
            "   “Her creatinine is 1.9” must come from her record — not from a guideline.",
        ],
    )

    add_slide(
        prs,
        "Part 3 — Supervisor + two workers",
        [
            "• Supervisor decides what to do and writes down why (inspectable).",
            "• Worker A — Intake extractor: reads the uploaded document.",
            "• Worker B — Evidence retriever: pulls matching guideline snippets.",
            "• Then we write one grounded answer for the doctor.",
            "• We keep the design small on purpose — two workers, not five.",
            "• A third “critic AI” is optional later; Week 1 verification already strips",
            "   claims that don’t have a source.",
        ],
    )

    add_slide(
        prs,
        "Part 4 — The quality gate (non-negotiable)",
        [
            "• 50 test cases that try to break us: bad scans, missing data, refusals,",
            "   “ignore the rules” prompts, and “don’t leak private data in logs.”",
            "• Scoring is yes/no (boolean) — not vague 1–10 stars.",
            "• Categories: valid schema, citations present, facts match sources,",
            "   safe refusal, no private data in logs.",
            "• If quality drops more than 5%, the build fails — demo can’t ship broken.",
            "• Graders will intentionally break something; our gate must catch it.",
        ],
        note="HARD GATE",
    )

    add_slide(
        prs,
        "Trust & safety (plain English)",
        [
            "• Only demo / fake patient data — no real patient information.",
            "• Logs keep IDs and timings — not names, not raw document text.",
            "• Permission check runs before any document is read.",
            "• Nurses/admins get the same rules as Week 1 (refuse when not allowed).",
            "• Doctor can click a citation and see the highlight on the PDF.",
        ],
    )

    add_slide(
        prs,
        "What we are NOT building this week",
        [
            "• Not five document types — only lab PDF + intake form until those work.",
            "• Not a giant pile of AI frameworks for show.",
            "• Not letting vision AI answer without filling the form first.",
            "• Not a mysterious supervisor — every handoff is logged.",
            "• Not guessing medical advice — retrieve, cite, check.",
        ],
        note="Narrower and stronger",
    )

    add_slide(
        prs,
        "Timeline & what this defense is",
        [
            "• Now — Architecture Defense (this plan).",
            "• Tuesday — MVP: two docs, two workers, RAG, 50-case gate.",
            "• Thursday — Early submission (harder / cleaner).",
            "• Sunday noon — Final: live app, demo video, cost & speed report.",
            "",
            "Today we are defending the plan — not shipping all the code yet.",
        ],
    )

    add_slide(
        prs,
        "Closing — one sentence",
        [
            "Week 2 = Week 1 trust",
            "        + eyes (read documents carefully)",
            "        + librarian (find guideline evidence)",
            "        + small crew (supervisor + two workers)",
            "        + bouncer (tests that fail the build if quality slips).",
            "",
            "Questions?",
        ],
    )

    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()
